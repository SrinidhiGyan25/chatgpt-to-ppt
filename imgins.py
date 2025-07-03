#!/usr/bin/env python3
"""
PowerPoint Image Inserter Script with Document Mapping and Collision Detection
Inserts images into PowerPoint slides with automatic position management and slide creation
"""

import os
import sys
import json
import csv
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Union, Set
from pptx import Presentation
from pptx.util import Inches, Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE
from dataclasses import dataclass
from enum import Enum

class Position(Enum):
    TOP_LEFT = "top-left"
    TOP_RIGHT = "top-right"
    BOTTOM_LEFT = "bottom-left"
    BOTTOM_RIGHT = "bottom-right"
    CENTER = "center"
    CUSTOM = "custom"

@dataclass
class ImagePlacement:
    image_number: int
    slide_number: int
    position: Position
    left: Optional[float] = None
    top: Optional[float] = None
    width: Optional[float] = None
    height: Optional[float] = None

class SlidePositionTracker:
    """Track occupied positions on slides"""
    
    def __init__(self):
        # Dictionary: slide_number -> set of occupied positions
        self.occupied_positions: Dict[int, Set[Position]] = {}
        # Available positions in order of preference
        self.available_positions = [
            Position.BOTTOM_LEFT,
            Position.BOTTOM_RIGHT,
            Position.TOP_RIGHT
        ]
    
    def is_position_occupied(self, slide_number: int, position: Position) -> bool:
        """Check if a position is occupied on a slide"""
        if slide_number not in self.occupied_positions:
            return False
        return position in self.occupied_positions[slide_number]
    
    def occupy_position(self, slide_number: int, position: Position):
        """Mark a position as occupied on a slide"""
        if slide_number not in self.occupied_positions:
            self.occupied_positions[slide_number] = set()
        self.occupied_positions[slide_number].add(position)
    
    def get_next_available_position(self, slide_number: int) -> Optional[Position]:
        """Get the next available position on a slide"""
        for position in self.available_positions:
            if not self.is_position_occupied(slide_number, position):
                return position
        return None
    
    def is_slide_full(self, slide_number: int) -> bool:
        """Check if all available positions on a slide are occupied"""
        return self.get_next_available_position(slide_number) is None
    
    def get_occupied_positions(self, slide_number: int) -> Set[Position]:
        """Get all occupied positions on a slide"""
        return self.occupied_positions.get(slide_number, set())

class PPTImageInserter:
    def __init__(self):
        self.image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'}
        self.image_files = []
        self.position_tracker = SlidePositionTracker()
        
    def get_image_files(self, directory: str = ".") -> List[Path]:
        """Get list of image files in directory"""
        image_files = []
        directory_path = Path(directory)
        
        if not directory_path.exists():
            raise FileNotFoundError(f"Image directory '{directory}' not found")
            
        for file_path in directory_path.glob("*"):
            if file_path.suffix.lower() in self.image_extensions:
                image_files.append(file_path)
        
        return sorted(image_files)
    
    def scan_existing_images(self, prs: Presentation):
        """Scan existing presentation for images and mark positions as occupied"""
        print("Scanning existing images in presentation...")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Determine position based on coordinates
                    position = self._determine_position_from_coordinates(
                        shape.left, shape.top, Inches(10), Inches(7.5)
                    )
                    self.position_tracker.occupy_position(slide_num, position)
                    print(f"  Slide {slide_num}: Found existing image at {position.value}")
    
    def _determine_position_from_coordinates(self, left, top, slide_width, slide_height) -> Position:
        """Determine position enum from coordinates"""
        # Convert to inches for comparison
        left_inches = left.inches if hasattr(left, 'inches') else left / 914400  # EMU to inches
        top_inches = top.inches if hasattr(top, 'inches') else top / 914400
        
        # Define thresholds (in inches)
        left_threshold = 3.0  # Center of 10-inch slide
        top_threshold = 2.5   # Center of 7.5-inch slide
        
        if left_inches < left_threshold and top_inches < top_threshold:
            return Position.TOP_LEFT
        elif left_inches >= left_threshold and top_inches < top_threshold:
            return Position.TOP_RIGHT
        elif left_inches < left_threshold and top_inches >= top_threshold:
            return Position.BOTTOM_LEFT
        elif left_inches >= left_threshold and top_inches >= top_threshold:
            return Position.BOTTOM_RIGHT
        else:
            return Position.CENTER
    
    def parse_mapping_document(self, mapping_file: str) -> List[ImagePlacement]:
        """Parse mapping document (supports JSON, CSV, or TXT formats)"""
        mapping_path = Path(mapping_file)
        
        if not mapping_path.exists():
            raise FileNotFoundError(f"Mapping file '{mapping_file}' not found")
        
        file_extension = mapping_path.suffix.lower()
        
        if file_extension == '.json':
            return self._parse_json_mapping(mapping_path)
        elif file_extension == '.csv':
            return self._parse_csv_mapping(mapping_path)
        elif file_extension in ['.txt', '.map']:
            return self._parse_txt_mapping(mapping_path)
        else:
            raise ValueError(f"Unsupported mapping file format: {file_extension}")
    
    def _parse_json_mapping(self, file_path: Path) -> List[ImagePlacement]:
        """Parse JSON mapping file"""
        with open(file_path, 'r') as f:
            data = json.load(f)
        
        placements = []
        for item in data:
            position_str = item.get('position', 'auto')
            position = Position(position_str) if position_str != 'auto' else None
            
            placement = ImagePlacement(
                image_number=item.get('image_number'),
                slide_number=item.get('slide_number', 'auto'),
                position=position,
                left=item.get('left'),
                top=item.get('top'),
                width=item.get('width'),
                height=item.get('height')
            )
            placements.append(placement)
        
        return placements
    
    def _parse_csv_mapping(self, file_path: Path) -> List[ImagePlacement]:
        """Parse CSV mapping file"""
        placements = []
        
        with open(file_path, 'r', newline='') as f:
            reader = csv.DictReader(f)
            for row in reader:
                position_str = row.get('position', 'auto')
                position = Position(position_str) if position_str != 'auto' else None
                
                slide_number = row.get('slide_number', 'auto')
                if slide_number != 'auto':
                    slide_number = int(slide_number) if slide_number else 'auto'
                
                placement = ImagePlacement(
                    image_number=int(row['image_number']) if row.get('image_number') else None,
                    slide_number=slide_number,
                    position=position,
                    left=float(row['left']) if row.get('left') and row['left'].strip() else None,
                    top=float(row['top']) if row.get('top') and row['top'].strip() else None,
                    width=float(row['width']) if row.get('width') and row['width'].strip() else None,
                    height=float(row['height']) if row.get('height') and row['height'].strip() else None
                )
                placements.append(placement)
        
        return placements
    
    def _parse_txt_mapping(self, file_path: Path) -> List[ImagePlacement]:
        """Parse TXT mapping file
        Format: image_number:slide_number:position:left:top:width:height
        Use 'auto' for automatic slide/position assignment
        """
        placements = []
        
        with open(file_path, 'r') as f:
            for line_num, line in enumerate(f, 1):
                line = line.strip()
                if not line or line.startswith('#'):
                    continue
                
                if ':' in line:
                    parts = line.split(':')
                else:
                    parts = line.split()
                
                if len(parts) < 1:
                    print(f"Warning: Skipping invalid line {line_num}: {line}")
                    continue
                
                try:
                    # Handle 'auto' values
                    slide_number = parts[1].strip() if len(parts) > 1 and parts[1].strip() else 'auto'
                    if slide_number != 'auto':
                        slide_number = int(slide_number)
                    
                    position_str = parts[2].strip() if len(parts) > 2 and parts[2].strip() else 'auto'
                    position = Position(position_str) if position_str != 'auto' else None
                    
                    placement = ImagePlacement(
                        image_number=int(parts[0]) if parts[0].strip() else None,
                        slide_number=slide_number,
                        position=position,
                        left=float(parts[3]) if len(parts) > 3 and parts[3].strip() else None,
                        top=float(parts[4]) if len(parts) > 4 and parts[4].strip() else None,
                        width=float(parts[5]) if len(parts) > 5 and parts[5].strip() else None,
                        height=float(parts[6]) if len(parts) > 6 and parts[6].strip() else None
                    )
                    placements.append(placement)
                except (ValueError, IndexError) as e:
                    print(f"Warning: Error parsing line {line_num}: {line} - {e}")
                    continue
        
        return placements
    
    def get_position_coordinates(self, position: Position, width=None, height=None) -> Tuple[float, float]:
        """Get coordinates for predefined positions"""
        positions = {
            Position.TOP_LEFT: (0.5, 0.5),
            Position.TOP_RIGHT: (6.0, 0.5),
            Position.BOTTOM_LEFT: (0.5, 5.0),
            Position.BOTTOM_RIGHT: (6.0, 5.0),
        }
        
        if position in positions:
            return positions[position]
        else:
            # For center, calculate dynamically
            slide_width = Inches(10)
            slide_height = Inches(7.5)
            img_width = Inches(width) if width else Inches(3)  # Default width
            img_height = Inches(height) if height else Inches(2)  # Default height
            
            left = (slide_width - img_width) / 2
            top = (slide_height - img_height) / 2
            return (left.inches, top.inches)
    
    def auto_assign_position(self, prs: Presentation, placement: ImagePlacement) -> Tuple[int, Position]:
        """Automatically assign slide and position for an image"""
        current_slide = 1
        
        # If slide is specified but position is auto
        if placement.slide_number != 'auto' and placement.position is None:
            slide_num = placement.slide_number
            if slide_num <= len(prs.slides):
                available_pos = self.position_tracker.get_next_available_position(slide_num)
                if available_pos:
                    return slide_num, available_pos
                else:
                    print(f"  Warning: Slide {slide_num} is full, creating new slide")
        
        # Find available position in existing slides
        for slide_num in range(1, len(prs.slides) + 1):
            available_pos = self.position_tracker.get_next_available_position(slide_num) 
            if available_pos:
                return slide_num, available_pos
        
        # All slides are full, create new slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        new_slide = prs.slides.add_slide(blank_layout)
        new_slide_num = len(prs.slides)
        print(f"  ✓ Created new slide {new_slide_num}")
        
        # Return first available position on new slide
        return new_slide_num, self.position_tracker.available_positions[0]
    
    def insert_image_to_slide(self, slide, image_path: Path, position: Position, 
                            left=None, top=None, width=None, height=None) -> bool:
        """Insert image into a slide at specified position"""
        try:
            if position == Position.CENTER:
                # Handle center positioning
                slide_width = Inches(10)
                slide_height = Inches(7.5)
                
                # Use provided dimensions or defaults
                img_width = Inches(width) if width else None
                img_height = Inches(height) if height else None
                
                if img_width and img_height:
                    # Calculate center position with specified dimensions
                    center_left = (slide_width - img_width) / 2
                    center_top = (slide_height - img_height) / 2
                    picture = slide.shapes.add_picture(str(image_path), center_left, center_top, 
                                                     img_width, img_height)
                else:
                    # Add with natural dimensions first, then center
                    temp_pic = slide.shapes.add_picture(str(image_path), Inches(0), Inches(0))
                    actual_width = temp_pic.width
                    actual_height = temp_pic.height
                    
                    # Remove temporary image
                    sp = temp_pic._element
                    sp.getparent().remove(sp)
                    
                    # Calculate center and add final image
                    center_left = (slide_width - actual_width) / 2
                    center_top = (slide_height - actual_height) / 2
                    picture = slide.shapes.add_picture(str(image_path), center_left, center_top, 
                                                     actual_width, actual_height)
            
            elif position == Position.CUSTOM:
                # Custom position with provided coordinates
                left_pos = Inches(left) if left else Inches(1)
                top_pos = Inches(top) if top else Inches(1)
                
                if width and height:
                    picture = slide.shapes.add_picture(str(image_path), left_pos, top_pos, 
                                                     Inches(width), Inches(height))
                else:
                    picture = slide.shapes.add_picture(str(image_path), left_pos, top_pos)
            
            else:
                # Predefined positions
                pos_left, pos_top = self.get_position_coordinates(position, width, height)
                left_pos = Inches(pos_left)
                top_pos = Inches(pos_top)
                
                if width and height:
                    picture = slide.shapes.add_picture(str(image_path), left_pos, top_pos,
                                                     Inches(width), Inches(height))
                else:
                    picture = slide.shapes.add_picture(str(image_path), left_pos, top_pos)
            
            return True
            
        except Exception as e:
            print(f"✗ Error inserting image {image_path.name}: {e}")
            return False
    
    def process_mappings(self, ppt_file: str, image_dir: str, mapping_file: str, output_file: str = None):
        """Main processing function with collision detection"""
        print("=== PowerPoint Image Inserter with Collision Detection ===\n")
        
        # Load or create presentation
        if os.path.exists(ppt_file):
            try:
                prs = Presentation(ppt_file)
                print(f"✓ Loaded presentation: {ppt_file}")
                # Scan for existing images
                self.scan_existing_images(prs)
            except Exception as e:
                print(f"✗ Error loading presentation: {e}")
                return False
        else:
            print(f"File '{ppt_file}' not found. Creating new presentation...")
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]
            prs.slides.add_slide(blank_slide_layout)
            print("✓ Created new presentation with 1 blank slide.")
        
        print(f"Presentation has {len(prs.slides)} slides.\n")
        
        # Get image files
        try:
            self.image_files = self.get_image_files(image_dir)
            print(f"✓ Found {len(self.image_files)} image files in '{image_dir}'")
            for i, img_file in enumerate(self.image_files, 1):
                print(f"  {i}. {img_file.name}")
            print()
        except Exception as e:
            print(f"✗ Error loading images: {e}")
            return False
        
        # Parse mapping document
        try:
            placements = self.parse_mapping_document(mapping_file)
            print(f"✓ Parsed {len(placements)} image placements from '{mapping_file}'\n")
        except Exception as e:
            print(f"✗ Error parsing mapping file: {e}")
            return False
        
        # Process each placement
        success_count = 0
        for i, placement in enumerate(placements, 1):
            print(f"Processing placement {i}/{len(placements)}:")
            
            # Validate image number
            if not placement.image_number or placement.image_number < 1 or placement.image_number > len(self.image_files):
                print(f"  ✗ Invalid image number: {placement.image_number}")
                continue
            
            image_file = self.image_files[placement.image_number - 1]
            print(f"  Image: {image_file.name}")
            
            # Handle automatic slide/position assignment
            if placement.slide_number == 'auto' or placement.position is None:
                assigned_slide, assigned_position = self.auto_assign_position(prs, placement)
                print(f"  Auto-assigned: Slide {assigned_slide}, Position {assigned_position.value}")
            else:
                assigned_slide = placement.slide_number
                assigned_position = placement.position
                
                # Check for collisions
                if self.position_tracker.is_position_occupied(assigned_slide, assigned_position):
                    print(f"  ⚠ Collision detected: Slide {assigned_slide}, Position {assigned_position.value} is occupied")
                    
                    # Try to find alternative
                    alt_position = self.position_tracker.get_next_available_position(assigned_slide)
                    if alt_position:
                        print(f"  ✓ Using alternative position: {alt_position.value}")
                        assigned_position = alt_position
                    else:
                        print(f"  ⚠ Slide {assigned_slide} is full, assigning to new slide")
                        assigned_slide, assigned_position = self.auto_assign_position(prs, placement)
                        print(f"  ✓ Assigned: Slide {assigned_slide}, Position {assigned_position.value}")
            
            # Validate final slide assignment
            if assigned_slide > len(prs.slides):
                print(f"  ✗ Invalid slide number: {assigned_slide}")
                continue
            
            slide = prs.slides[assigned_slide - 1]
            
            # Insert image
            success = self.insert_image_to_slide(
                slide, image_file, assigned_position,
                placement.left, placement.top, placement.width, placement.height
            )
            
            if success:
                # Mark position as occupied
                self.position_tracker.occupy_position(assigned_slide, assigned_position)
                print(f"  ✓ Successfully inserted into Slide {assigned_slide} at {assigned_position.value}")
                success_count += 1
            else:
                print(f"  ✗ Failed to insert image")
            
            print()  # Add spacing between placements
        
        # Save presentation
        try:
            if not output_file:
                output_file = ppt_file
            prs.save(output_file)
            print(f"✓ Successfully processed {success_count}/{len(placements)} placements")
            print(f"✓ Final presentation has {len(prs.slides)} slides")
            print(f"✓ Presentation saved as: {output_file}")
            
            # Show final position summary
            self._print_position_summary()
            return True
        except Exception as e:
            print(f"\n✗ Error saving presentation: {e}")
            return False
    
    def _print_position_summary(self):
        """Print summary of occupied positions"""
        print("\n=== Position Summary ===")
        for slide_num in sorted(self.position_tracker.occupied_positions.keys()):
            occupied = self.position_tracker.get_occupied_positions(slide_num)
            occupied_list = [pos.value for pos in occupied]
            print(f"Slide {slide_num}: {', '.join(occupied_list)}")

def create_sample_mapping_files():
    """Create sample mapping files with auto-assignment examples"""
    
    # Sample JSON mapping with auto assignment
    json_sample = [
        {"image_number": 1, "slide_number": "auto", "position": "auto"},
        {"image_number": 2, "slide_number": "auto", "position": "auto"},
        {"image_number": 3, "slide_number": 1, "position": "bottom-left"},
        {"image_number": 4, "slide_number": "auto", "position": "auto"},
        {"image_number": 5, "slide_number": "auto", "position": "top-right"}
    ]
    
    with open('sample_mapping.json', 'w') as f:
        json.dump(json_sample, f, indent=2)
    
    # Sample CSV mapping
    csv_content = """image_number,slide_number,position,left,top,width,height
1,auto,auto,,,3.0,2.0
2,auto,auto,,,3.0,2.0
3,1,bottom-left,,,3.0,2.0
4,auto,auto,,,3.0,2.0
5,auto,top-right,,,3.0,2.0"""
    
    with open('sample_mapping.csv', 'w') as f:
        f.write(csv_content)
    
    # Sample TXT mapping
    txt_content = """# Image mapping file with auto-assignment
# Format: image_number:slide_number:position:left:top:width:height
# Use 'auto' for automatic assignment
1:auto:auto:::3.0:2.0
2:auto:auto:::3.0:2.0
3:1:bottom-left:::3.0:2.0
4:auto:auto:::3.0:2.0
5:auto:top-right:::3.0:2.0"""
    
    with open('sample_mapping.txt', 'w') as f:
        f.write(txt_content)
    
    print("✓ Created sample mapping files with auto-assignment examples")
    print("Available positions: bottom-left, bottom-right, top-right")
    print("Use 'auto' for automatic slide/position assignment")

def main():
    if len(sys.argv) > 1 and sys.argv[1] == '--create-samples':
        create_sample_mapping_files()
        return
    
    inserter = PPTImageInserter()
    
    try:
        # Get inputs
        ppt_file = input("Enter PowerPoint file path (or press Enter for 'presentation.pptx'): ").strip()
        if not ppt_file:
            ppt_file = "presentation.pptx"
        
        image_dir = input("Enter images directory (press Enter for current directory): ").strip()
        if not image_dir:
            image_dir = "."
        
        mapping_file = input("Enter mapping file path (JSON/CSV/TXT): ").strip()
        if not mapping_file:
            print("Mapping file is required!")
            return
        
        output_file = input(f"Save as (press Enter for '{ppt_file}'): ").strip()
        if not output_file:
            output_file = ppt_file
        
        # Process mappings with collision detection
        inserter.process_mappings(ppt_file, image_dir, mapping_file, output_file)
        
    except KeyboardInterrupt:
        print("\n\nOperation cancelled by user.")
    except Exception as e:
        print(f"\nUnexpected error: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()