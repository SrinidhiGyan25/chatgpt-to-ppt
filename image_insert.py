#!/usr/bin/env python3
"""
PowerPoint Image Inserter Script
Inserts images into PowerPoint slides based on user input
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE

def get_slide_info(prs):
    """Display information about existing slides"""
    print(f"\nPresentation has {len(prs.slides)} slides:")
    for i, slide in enumerate(prs.slides, 1):
        layout_name = slide.slide_layout.name
        print(f"  Slide {i}: Layout '{layout_name}'")
    print()

def get_image_files(directory="."):
    """Get list of image files in directory"""
    image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'}
    image_files = []
    
    for file_path in Path(directory).glob("*"):
        if file_path.suffix.lower() in image_extensions:
            image_files.append(file_path)
    
    return sorted(image_files)

def display_images(image_files):
    """Display available image files"""
    if not image_files:
        print("No image files found in current directory.")
        return False
    
    print("Available image files:")
    for i, img_file in enumerate(image_files, 1):
        print(f"  {i}. {img_file.name}")
    print()
    return True

def get_position_input():
    """Get position and size input from user"""
    print("Enter image position and size:")
    print("Options:")
    print("1. Center (default size)")
    print("2. Top-left corner")
    print("3. Top-right corner") 
    print("4. Bottom-left corner")
    print("5. Bottom-right corner")
    print("6. Custom position")
    
    choice = input("Choose position (1-6) or press Enter for center: ").strip()
    
    if choice == "1" or choice == "":
        return "center", None, None, None, None
    elif choice == "2":
        return "custom", Inches(0.5), Inches(0.5), None, None
    elif choice == "3":
        return "custom", Inches(6), Inches(0.5), None, None
    elif choice == "4":
        return "custom", Inches(0.5), Inches(5), None, None
    elif choice == "5":
        return "custom", Inches(6), Inches(5), None, None
    elif choice == "6":
        try:
            left = float(input("Left position (inches): "))
            top = float(input("Top position (inches): "))
            width_input = input("Width (inches, press Enter for auto): ").strip()
            height_input = input("Height (inches, press Enter for auto): ").strip()
            
            width = Inches(float(width_input)) if width_input else None
            height = Inches(float(height_input)) if height_input else None
            
            return "custom", Inches(left), Inches(top), width, height
        except ValueError:
            print("Invalid input. Using center position.")
            return "center", None, None, None, None
    else:
        print("Invalid choice. Using center position.")
        return "center", None, None, None, None

def insert_image_to_slide(slide, image_path, position_type, left=None, top=None, width=None, height=None):
    """Insert image into a slide at specified position"""
    try:
        if position_type == "center":
            # Calculate center position
            slide_width = Inches(10)  # Standard slide width
            slide_height = Inches(7.5)  # Standard slide height
            
            # Add image first to get its natural dimensions
            temp_pic = slide.shapes.add_picture(str(image_path), Inches(0), Inches(0))
            
            # Calculate centered position
            img_width = temp_pic.width
            img_height = temp_pic.height
            
            # Remove temporary image
            sp = temp_pic._element
            sp.getparent().remove(sp)
            
            # Calculate center position
            left = (slide_width - img_width) / 2
            top = (slide_height - img_height) / 2
            
            # Add image at center
            picture = slide.shapes.add_picture(str(image_path), left, top, img_width, img_height)
        else:
            # Custom position
            if width and height:
                picture = slide.shapes.add_picture(str(image_path), left, top, width, height)
            elif width:
                picture = slide.shapes.add_picture(str(image_path), left, top, width=width)
            elif height:
                picture = slide.shapes.add_picture(str(image_path), left, top, height=height)
            else:
                picture = slide.shapes.add_picture(str(image_path), left, top)
        
        print(f"✓ Image inserted successfully")
        return True
        
    except Exception as e:
        print(f"✗ Error inserting image: {e}")
        return False

def main():
    print("=== PowerPoint Image Inserter ===\n")
    
    # Get PowerPoint file
    ppt_file = input("Enter PowerPoint file path (or press Enter for 'presentation.pptx'): ").strip()
    if not ppt_file:
        ppt_file = "presentation.pptx"
    
    # Check if file exists
    if not os.path.exists(ppt_file):
        print(f"File '{ppt_file}' not found.")
        create_new = input("Create new presentation? (y/n): ").lower().startswith('y')
        if create_new:
            prs = Presentation()
            # Add a blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            prs.slides.add_slide(blank_slide_layout)
            print("Created new presentation with 1 blank slide.")
        else:
            return
    else:
        try:
            prs = Presentation(ppt_file)
            print(f"Loaded presentation: {ppt_file}")
        except Exception as e:
            print(f"Error loading presentation: {e}")
            return
    
    # Display slide information
    get_slide_info(prs)
    
    # Get image directory
    img_dir = input("Enter image directory (press Enter for current directory): ").strip()
    if not img_dir:
        img_dir = "."
    
    # Get available images
    image_files = get_image_files(img_dir)
    if not display_images(image_files):
        return
    
    # Main loop for inserting images
    while True:
        print("\n" + "="*50)
        
        # Select slide
        try:
            slide_input = input(f"Enter slide number (1-{len(prs.slides)}) or 'q' to quit: ").strip()
            if slide_input.lower() == 'q':
                break
                
            slide_num = int(slide_input)
            if slide_num < 1 or slide_num > len(prs.slides):
                print("Invalid slide number.")
                continue
                
            slide = prs.slides[slide_num - 1]
            print(f"Selected slide {slide_num}")
            
        except ValueError:
            print("Invalid input. Please enter a number.")
            continue
        
        # Select image
        try:
            img_input = input(f"Enter image number (1-{len(image_files)}) or 's' to skip: ").strip()
            if img_input.lower() == 's':
                continue
                
            img_num = int(img_input)
            if img_num < 1 or img_num > len(image_files):
                print("Invalid image number.")
                continue
                
            selected_image = image_files[img_num - 1]
            print(f"Selected image: {selected_image.name}")
            
        except ValueError:
            print("Invalid input. Please enter a number.")
            continue
        
        # Get position
        position_type, left, top, width, height = get_position_input()
        
        # Insert image
        success = insert_image_to_slide(slide, selected_image, position_type, left, top, width, height)
        
        if success:
            # Ask if user wants to continue
            continue_choice = input("\nAdd another image? (y/n): ").lower()
            if not continue_choice.startswith('y'):
                break
    
    # Save presentation
    try:
        output_file = input(f"\nSave as (press Enter for '{ppt_file}'): ").strip()
        if not output_file:
            output_file = ppt_file
            
        prs.save(output_file)
        print(f"✓ Presentation saved as: {output_file}")
        
    except Exception as e:
        print(f"✗ Error saving presentation: {e}")

if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\n\nOperation cancelled by user.")
    except Exception as e:
        print(f"\nUnexpected error: {e}")
        sys.exit(1)